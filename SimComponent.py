# ============================================================
# EVMS Deck per Program – 3 Slides
#   1) EVMS Trend Overview
#   2) EVMS Detail – Sub Team Labor & Manpower
#   3) EVMS Detail – Sub Team CPI/SPI/BEI Metrics
# ============================================================

import os
from datetime import datetime
import numpy as np
import pandas as pd
import plotly.graph_objects as go

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ---------------- CONFIG ----------------

cobra_files = {
    "Abrams_STS_2022": "data/Cobra-Abrams STS 2022.xlsx",
    "Abrams_STS"     : "data/Cobra-Abrams STS.xlsx",
    "XM30"           : "data/Cobra-XM30.xlsx",
}

PENSKE_PATH = "data/OpenPlan_Activity-Penske.xlsx"
THEME_PATH  = "data/Theme.pptx"      # optional theme
OUTPUT_DIR  = "EVMS_Output"
os.makedirs(OUTPUT_DIR, exist_ok=True)

# Accounting calendar closings (2025 – Detroit area)
ACCOUNTING_CLOSINGS = {
    (2025, 1): 26,
    (2025, 2): 23,
    (2025, 3): 30,
    (2025, 4): 27,
    (2025, 5): 25,
    (2025, 6): 29,
    (2025, 7): 27,
    (2025, 8): 24,
    (2025, 9): 28,
    (2025,10): 26,
    (2025,11): 23,
    (2025,12): 31,
}

HOURS_PER_FTE = 173.33  # approx 9/80 schedule monthly hours

# Threshold colors from dashboard key
COLOR_BLUE_LIGHT = RGBColor(142, 180, 227)
COLOR_GREEN      = RGBColor( 51, 153, 102)
COLOR_YELLOW     = RGBColor(255, 255, 153)
COLOR_RED        = RGBColor(192,  80,  77)

# ============================================================
# SHARED HELPERS
# ============================================================

def normalize_columns(df: pd.DataFrame) -> pd.DataFrame:
    """Uppercase and strip spaces/hyphens/underscores from column names."""
    return df.rename(columns={
        c: c.strip().upper().replace(" ", "").replace("-", "").replace("_", "")
        for c in df.columns
    })

def map_cost_sets(cost_cols):
    """Map Cobra COST-SET labels to BCWS / BCWP / ACWP / ETC using fuzzy matching."""
    cleaned = {
        col: col.replace(" ", "").replace("-", "").replace("_", "").upper()
        for col in cost_cols
    }
    bcws = bcwp = acwp = etc = None
    for orig, clean in cleaned.items():
        if ("ACWP" in clean) or ("ACTUAL" in clean) or ("ACWPHRS" in clean):
            acwp = orig
        elif ("BCWS" in clean) or ("BUDGET" in clean) or ("PLAN" in clean):
            bcws = orig
        elif ("BCWP" in clean) or ("EARNED" in clean) or ("PROGRESS" in clean):
            bcwp = orig
        elif "ETC" in clean:
            etc = orig
    return bcws, bcwp, acwp, etc

def spi_cpi_color(x: float):
    """SPI / CPI / BEI thresholds."""
    if pd.isna(x):
        return None
    if x >= 1.05:
        return COLOR_BLUE_LIGHT
    elif x >= 0.98:
        return COLOR_GREEN
    elif x >= 0.95:
        return COLOR_YELLOW
    else:
        return COLOR_RED

def vac_color_from_ratio(r: float):
    """VAC/BAC thresholds."""
    if pd.isna(r):
        return None
    # ≥ +5% blue; +5%..−2% green; −2%..−5% yellow; < −5% red
    if r >= 0.05:
        return COLOR_BLUE_LIGHT
    elif r >= -0.02:
        return COLOR_GREEN
    elif r >= -0.05:
        return COLOR_YELLOW
    else:
        return COLOR_RED

def manpower_var_color(r: float):
    """Program manpower %Var thresholds based on Actual/Demand ratio."""
    if pd.isna(r):
        return None
    # ≥110% red; 110–105 yellow; 105–90 green; 90–85 yellow; <85 red
    if r >= 1.10:
        return COLOR_RED
    elif r >= 1.05:
        return COLOR_YELLOW
    elif r >= 0.90:
        return COLOR_GREEN
    elif r >= 0.85:
        return COLOR_YELLOW
    else:
        return COLOR_RED

def get_status_dates(dates: pd.Series):
    """Status dates from accounting calendar; fallback to last two EV dates."""
    dates = pd.to_datetime(dates)
    max_date = dates.max()

    closing_dates = []
    for (year, month), day in ACCOUNTING_CLOSINGS.items():
        d = datetime(year, month, day)
        if d <= max_date:
            closing_dates.append(d)
    closing_dates = sorted(closing_dates)

    if len(closing_dates) >= 2:
        curr = closing_dates[-1]
        prev = closing_dates[-2]
    elif len(closing_dates) == 1:
        curr = prev = closing_dates[0]
    else:
        uniq = sorted(dates.unique())
        curr = uniq[-1]
        prev = uniq[-2] if len(uniq) > 1 else uniq[-1]
    return curr, prev

def get_row_on_or_before(evdf: pd.DataFrame, date: datetime):
    sub = evdf[evdf["DATE"] <= date]
    if sub.empty:
        return evdf.iloc[0]
    return sub.iloc[-1]

# ============================================================
# EVMS FROM COBRA
# ============================================================

def compute_ev_timeseries(df: pd.DataFrame) -> pd.DataFrame:
    """Build Monthly & Cumulative CPI/SPI from Cobra data."""
    df = df.copy()
    df["DATE"] = pd.to_datetime(df["DATE"])

    pivot = df.pivot_table(
        index="DATE", columns="COSTSET", values="HOURS", aggfunc="sum"
    ).reset_index()

    cost_cols = [c for c in pivot.columns if c != "DATE"]
    bcws_col, bcwp_col, acwp_col, _ = map_cost_sets(cost_cols)

    missing = []
    if bcws_col is None: missing.append("BCWS")
    if bcwp_col is None: missing.append("BCWP")
    if acwp_col is None: missing.append("ACWP")
    if missing:
        raise ValueError(f"Missing cost sets {missing}; found {cost_cols}")

    BCWS_raw = pivot[bcws_col].fillna(0.0)
    BCWP_raw = pivot[bcwp_col].fillna(0.0)
    ACWP_raw = pivot[acwp_col].fillna(0.0)

    monthly_cpi = BCWP_raw / ACWP_raw.replace(0, np.nan)
    monthly_spi = BCWP_raw / BCWS_raw.replace(0, np.nan)

    cum_bcws = BCWS_raw.cumsum()
    cum_bcwp = BCWP_raw.cumsum()
    cum_acwp = ACWP_raw.cumsum()

    cumulative_cpi = cum_bcwp / cum_acwp.replace(0, np.nan)
    cumulative_spi = cum_bcwp / cum_bcws.replace(0, np.nan)

    return pd.DataFrame({
        "DATE": pivot["DATE"],
        "Monthly CPI": monthly_cpi,
        "Monthly SPI": monthly_spi,
        "Cumulative CPI": cumulative_cpi,
        "Cumulative SPI": cumulative_spi,
    })

# ============================================================
# BEI FROM OPENPLAN PENSKE
# ============================================================

PENSKE = pd.read_excel(PENSKE_PATH)
PENSKE = normalize_columns(PENSKE)

def _bei_for_df(df: pd.DataFrame, base_col: str, act_col: str, as_of: datetime):
    denom = df[df[base_col] <= as_of]
    if denom.empty:
        return np.nan
    numer = denom[denom[act_col].notna() & (denom[act_col] <= as_of)]
    return len(numer) / len(denom)

def compute_program_bei(program_name: str,
                        status_date: datetime,
                        prev_status_date: datetime):
    """Program-level BEI CTD/LSD from Penske."""
    df = PENSKE.copy()

    if "PROGRAM" in df.columns:
        mask = df["PROGRAM"].astype(str).str.contains(program_name, case=False, na=False)
        if mask.any():
            df = df[mask]

    base_col = next((c for c in df.columns if "BASELINEFINISH" in c or ("BASELINE" in c and "FINISH" in c)), None)
    act_col  = next((c for c in df.columns if "ACTUALFINISH"   in c or ("ACTUAL"   in c and "FINISH" in c)), None)
    if base_col is None or act_col is None:
        return np.nan, np.nan, None, None

    df[base_col] = pd.to_datetime(df[base_col], errors="coerce")
    df[act_col]  = pd.to_datetime(df[act_col],  errors="coerce")

    lev_col = next((c for c in df.columns if "LEVTYPE" in c), None)
    if lev_col is not None:
        df = df[~df[lev_col].isin(["A", "B"])]   # drop LOE & milestones

    if df.empty:
        return np.nan, np.nan, None, None

    bei_ctd = _bei_for_df(df, base_col, act_col, status_date)
    bei_lsd = _bei_for_df(df, base_col, act_col, prev_status_date)
    return bei_ctd, bei_lsd, df, (base_col, act_col)

def compute_bei_by_subteam(program_name: str,
                           subteams,
                           status_date: datetime,
                           prev_status_date: datetime):
    """
    Return dict: subteam -> (BEI_CTD, BEI_LSD)
    Uses TEAM for mapping when available; falls back to program-level BEI.
    """
    prog_ctd, prog_lsd, df_prog, cols = compute_program_bei(
        program_name, status_date, prev_status_date
    )
    if df_prog is None or cols is None:
        return {st: (np.nan, np.nan) for st in subteams}

    base_col, act_col = cols
    team_col = "TEAM" if "TEAM" in df_prog.columns else None

    result = {}
    for st in subteams:
        if team_col:
            mask = df_prog[team_col].astype(str).str.upper() == str(st).upper()
            sub = df_prog[mask]
        else:
            sub = pd.DataFrame()

        if sub.empty:
            result[st] = (prog_ctd, prog_lsd)
        else:
            bei_ctd = _bei_for_df(sub, base_col, act_col, status_date)
            bei_lsd = _bei_for_df(sub, base_col, act_col, prev_status_date)
            result[st] = (bei_ctd, bei_lsd)
    return result

# ============================================================
# TABLE BUILDERS
# ============================================================

def build_labor_table(cobra: pd.DataFrame) -> pd.DataFrame:
    """
    Labor Hours Performance by Sub Team:
      %COMP = CTD BCWP / CTD BAC
      BAC   = CTD BCWS
      EAC   = CTD ACWP + CTD ETC
      VAC   = BAC - EAC
    """
    df = cobra.copy()
    if "SUBTEAM" not in df.columns:
        raise ValueError("Labor table requires SUBTEAM column (from SUB_TEAM).")

    pivot = df.pivot_table(
        index="SUBTEAM", columns="COSTSET", values="HOURS", aggfunc="sum"
    ).fillna(0)

    bcws_col, bcwp_col, acwp_col, etc_col = map_cost_sets(pivot.columns)

    missing = []
    if bcws_col is None: missing.append("BCWS")
    if bcwp_col is None: missing.append("BCWP")
    if acwp_col is None: missing.append("ACWP")
    if missing:
        raise ValueError(f"Missing cost sets for labor table: {missing}")

    BAC   = pivot[bcws_col]
    BCWP  = pivot[bcwp_col]
    ACWP  = pivot[acwp_col]
    ETC   = pivot[etc_col] if (etc_col is not None and etc_col in pivot.columns) else 0.0
    EAC   = ACWP + ETC
    VAC   = BAC - EAC
    pct_c = np.where(BAC != 0, BCWP / BAC, np.nan)

    labor_df = pd.DataFrame({
        "Sub Team": pivot.index,
        "%COMP": pct_c,
        "BAC": BAC,
        "EAC": EAC,
        "VAC": VAC,
    }).reset_index(drop=True)

    return labor_df

def build_subteam_metric_table(cobra: pd.DataFrame,
                               program_name: str,
                               curr_date: datetime,
                               prev_date: datetime) -> pd.DataFrame:
    """
    Sub-team level EVMS metrics:
      CPI CTD/LSD, SPI CTD/LSD, BEI CTD/LSD
    """
    df = cobra.copy()
    if "SUBTEAM" not in df.columns:
        raise ValueError("Metric table requires SUBTEAM column (from SUB_TEAM).")

    subteams = sorted(df["SUBTEAM"].unique())
    bei_map = compute_bei_by_subteam(program_name, subteams, curr_date, prev_date)

    out_rows = []
    for st in subteams:
        sub_df = df[df["SUBTEAM"] == st]
        try:
            ev_sub = compute_ev_timeseries(sub_df)
        except ValueError:
            ev_sub = None

        if ev_sub is not None and not ev_sub.empty:
            row_curr = get_row_on_or_before(ev_sub, curr_date)
            row_prev = get_row_on_or_before(ev_sub, prev_date)
            cpi_ctd = row_curr["Cumulative CPI"]
            cpi_lsd = row_prev["Cumulative CPI"]
            spi_ctd = row_curr["Cumulative SPI"]
            spi_lsd = row_prev["Cumulative SPI"]
        else:
            cpi_ctd = cpi_lsd = spi_ctd = spi_lsd = np.nan

        bei_ctd, bei_lsd = bei_map.get(st, (np.nan, np.nan))

        out_rows.append({
            "Sub Team": st,
            "CPI CTD": cpi_ctd,
            "CPI LSD": cpi_lsd,
            "SPI CTD": spi_ctd,
            "SPI LSD": spi_lsd,
            "BEI CTD": bei_ctd,
            "BEI LSD": bei_lsd,
        })

    metrics_df = pd.DataFrame(out_rows)
    return metrics_df

def build_manpower_table(cobra: pd.DataFrame,
                         curr_date: datetime,
                         prev_date: datetime) -> pd.DataFrame:
    """
    Program Manpower (FTE) – CTD based:
      Demand FTE (CTD), Actual FTE (CTD), % Var,
      Next Mo BCWS FTE, Next Mo ETC FTE
    """
    df = cobra.copy()
    df["DATE"] = pd.to_datetime(df["DATE"])

    pivot = df.pivot_table(
        index=df["DATE"].dt.to_period("M"),
        columns="COSTSET",
        values="HOURS",
        aggfunc="sum"
    ).fillna(0)

    if pivot.empty:
        return pd.DataFrame(columns=["Demand", "Actual", "% Var",
                                     "Next Mo BCWS", "Next Mo ETC"])

    bcws_col, _, acwp_col, etc_col = map_cost_sets(pivot.columns)
    if bcws_col is None or acwp_col is None:
        return pd.DataFrame(columns=["Demand", "Actual", "% Var",
                                     "Next Mo BCWS", "Next Mo ETC"])

    demand_hrs = pivot[bcws_col]
    actual_hrs = pivot[acwp_col]
    etc_hrs    = pivot[etc_col] if (etc_col is not None and etc_col in pivot.columns) else pd.Series(0, index=pivot.index)

    # CTD FTE at status month
    cum_demand_fte = demand_hrs.cumsum() / HOURS_PER_FTE
    cum_actual_fte = actual_hrs.cumsum() / HOURS_PER_FTE

    status_period = pd.Period(curr_date, freq="M")
    next_period   = status_period + 1

    demand = cum_demand_fte.get(status_period, np.nan)
    actual = cum_actual_fte.get(status_period, np.nan)

    if demand is not None and not pd.isna(demand) and demand != 0:
        pct_var = actual / demand
    else:
        pct_var = np.nan

    # Next month (non-CTD) BCWS & ETC in FTE
    next_bcws_fte = (demand_hrs / HOURS_PER_FTE).get(next_period, np.nan)
    next_etc_fte  = (etc_hrs    / HOURS_PER_FTE).get(next_period, np.nan)

    manpower_df = pd.DataFrame({
        "Demand": [demand],
        "Actual": [actual],
        "% Var": [pct_var],
        "Next Mo BCWS": [next_bcws_fte],
        "Next Mo ETC": [next_etc_fte],
    })
    return manpower_df

# ============================================================
# PLOTTING & PPT HELPERS
# ============================================================

def create_evms_plot(program: str, evdf: pd.DataFrame) -> go.Figure:
    fig = go.Figure()

    # EV performance background zones
    fig.add_hrect(y0=0.90, y1=0.95, fillcolor="red",       opacity=0.25, line_width=0)
    fig.add_hrect(y0=0.95, y1=0.98, fillcolor="yellow",    opacity=0.25, line_width=0)
    fig.add_hrect(y0=0.98, y1=1.05, fillcolor="green",     opacity=0.25, line_width=0)
    fig.add_hrect(y0=1.05, y1=1.20, fillcolor="lightblue", opacity=0.25, line_width=0)

    fig.add_trace(go.Scatter(
        x=evdf["DATE"], y=evdf["Monthly CPI"],
        mode="markers", name="Monthly CPI",
        marker=dict(symbol="diamond", size=7, color="yellow")
    ))
    fig.add_trace(go.Scatter(
        x=evdf["DATE"], y=evdf["Monthly SPI"],
        mode="markers", name="Monthly SPI",
        marker=dict(size=6, color="black")
    ))
    fig.add_trace(go.Scatter(
        x=evdf["DATE"], y=evdf["Cumulative CPI"],
        mode="lines", name="Cumulative CPI",
        line=dict(color="blue", width=3)
    ))
    fig.add_trace(go.Scatter(
        x=evdf["DATE"], y=evdf["Cumulative SPI"],
        mode="lines", name="Cumulative SPI",
        line=dict(color="darkgrey", width=3)
    ))

    fig.update_layout(
        xaxis_title="Month",
        yaxis_title="EV Index",
        yaxis=dict(range=[0.90, 1.20]),
        template="simple_white",
        height=360,
        margin=dict(l=60, r=20, t=20, b=50)
    )
    return fig

def get_blank_layout(prs: Presentation):
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower():
            return layout
    return prs.slides[0].slide_layout

def add_title(slide, text, left=0.5, top=0.3):
    box = slide.shapes.add_textbox(Inches(left), Inches(top),
                                   Inches(9.0), Inches(0.6))
    tf = box.text_frame
    tf.text = text
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(24)

def add_simple_table(slide, df: pd.DataFrame,
                     left_in, top_in, width_in, height_in=None):
    """
    Create a table for a DataFrame; return (pptx table object, height_inches).
    CPI/SPI/BEI metrics and CTD/LSD columns shown with 2 decimals.
    Percents (%COMP, % Var) shown with 2 decimal places.
    """
    rows, cols = df.shape
    base_h = 0.45
    row_h  = 0.24   # tighter rows
    if height_in is None:
        height_in = base_h + row_h * (rows + 1)

    shape = slide.shapes.add_table(
        rows + 1, cols,
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in)
    )
    table = shape.table

    # header
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)

    metric_cols = {"CPI CTD","CPI LSD","SPI CTD","SPI LSD","BEI CTD","BEI LSD","CTD","LSD"}

    # body
    for i in range(rows):
        for j, col in enumerate(df.columns):
            val = df.iloc[i, j]
            cell = table.cell(i + 1, j)

            if isinstance(val, (float, np.floating)):
                ucol = col.upper()
                if "%COMP" in ucol or "VAR" in ucol:
                    # percentages (ratio 0–1)
                    cell.text = "" if pd.isna(val) else f"{val*100:.2f}%"
                elif col in metric_cols:
                    cell.text = "" if pd.isna(val) else f"{val:.2f}"
                else:
                    cell.text = "" if pd.isna(val) else f"{val:,.1f}"
            else:
                cell.text = "" if pd.isna(val) else str(val)

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(9.5)

    return table, height_in

def add_rcca_box(slide, label="Comments / Root Cause & Corrective Actions",
                 left_in=0.5, top_in=5.3, width_in=9.0, height_in=1.0):
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in)
    )
    tf = box.text_frame
    tf.text = label + ":"
    p = tf.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(14)
    tf.add_paragraph()  # blank line for notes

def add_threshold_legend(slide, left_in=6.2, top_in=3.7, width_in=3.5, height_in=1.8):
    """Add a small legend describing the color thresholds."""
    box = slide.shapes.add_textbox(
        Inches(left_in), Inches(top_in),
        Inches(width_in), Inches(height_in)
    )
    tf = box.text_frame
    tf.text = "Color Coding Thresholds:"
    p0 = tf.paragraphs[0]
    p0.font.bold = True
    p0.font.size = Pt(10)

    lines = [
        "• CPI / SPI / BEI: Blue ≥ 1.05; Green 0.98–1.05; Yellow 0.95–0.98; Red < 0.95.",
        "• VAC/BAC: Blue ≥ +5%; Green +5% to −2%; Yellow −2% to −5%; Red < −5%.",
        "• Manpower % Var: Green 90–105%; Yellow 85–90% or 105–110%; Red <85% or ≥110%.",
    ]
    for text in lines:
        p = tf.add_paragraph()
        p.text = text
        p.level = 1
        p.font.size = Pt(8)

# ============================================================
# MAIN LOOP – Build decks
# ============================================================

for program, path in cobra_files.items():
    print(f"Processing {program} ...")

    cobra = pd.read_excel(path)
    cobra = normalize_columns(cobra)

    required = {"DATE", "COSTSET", "HOURS"}
    missing = required - set(cobra.columns)
    if missing:
        raise ValueError(f"{program}: missing columns {missing} after normalization.")

    # EV series and status dates
    evdf = compute_ev_timeseries(cobra)
    curr_date, prev_date = get_status_dates(evdf["DATE"])

    # Program-level CPI/SPI CTD & LSD
    row_curr = get_row_on_or_before(evdf, curr_date)
    row_prev = get_row_on_or_before(evdf, prev_date)
    cpi_ctd = row_curr["Cumulative CPI"]
    spi_ctd = row_curr["Cumulative SPI"]
    cpi_lsd = row_prev["Cumulative CPI"]
    spi_lsd = row_prev["Cumulative SPI"]

    # Program-level BEI CTD/LSD
    prog_bei_ctd, prog_bei_lsd, _, _ = compute_program_bei(program, curr_date, prev_date)

    metrics_df = pd.DataFrame({
        "Metric": ["CPI", "SPI", "BEI"],
        "CTD":    [cpi_ctd, spi_ctd, prog_bei_ctd],
        "LSD":    [cpi_lsd, spi_lsd, prog_bei_lsd],
    })

    print(f"  CTD date: {curr_date.date()}, LSD date: {prev_date.date()}")
    print(metrics_df.round(2), "\n")

    # EVMS plot
    fig = create_evms_plot(program, evdf)
    img_path = os.path.join(OUTPUT_DIR, f"{program}_EVMS.png")
    fig.write_image(img_path, scale=3)

    # Subteam tables & manpower
    labor_df    = build_labor_table(cobra)
    metrics_sub = build_subteam_metric_table(cobra, program, curr_date, prev_date)
    manpower_df = build_manpower_table(cobra, curr_date, prev_date)

    # --- PowerPoint deck per program ---
    prs = Presentation(THEME_PATH) if os.path.exists(THEME_PATH) else Presentation()
    blank_layout = get_blank_layout(prs)

    # ---------------------------------------------------
    # Slide 1 – EVMS Trend Overview
    # ---------------------------------------------------
    slide1 = prs.slides.add_slide(blank_layout)
    add_title(slide1, f"{program} EVMS Trend Overview")

    slide1.shapes.add_picture(
        img_path,
        Inches(0.5), Inches(1.0),
        width=Inches(5.5)
    )

    tbl1, _ = add_simple_table(slide1, metrics_df,
                               left_in=6.2, top_in=1.0,
                               width_in=3.5, height_in=1.2)

    # Color CTD/LSD cells
    for i in range(1, len(metrics_df) + 1):
        for col_name, col_idx in (("CTD", 1), ("LSD", 2)):
            val = metrics_df.iloc[i-1][col_name]
            rgb = spi_cpi_color(val)
            if rgb is not None:
                cell = tbl1.cell(i, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb

    add_threshold_legend(slide1)
    add_rcca_box(slide1)

    # ---------------------------------------------------
    # Slide 2 – Sub Team Labor & Program Manpower
    # ---------------------------------------------------
    slide2 = prs.slides.add_slide(blank_layout)
    add_title(slide2, f"{program} EVMS Detail – Sub Team Labor & Manpower")

    # Labor Hours Performance table
    top = 0.9
    labor_tbl, labor_h = add_simple_table(
        slide2, labor_df,
        left_in=0.5, top_in=top,
        width_in=9.0
    )

    # Color VAC cells by VAC/BAC
    cols_labor = list(labor_df.columns)
    vac_idx = cols_labor.index("VAC")
    for i in range(len(labor_df)):
        vac = labor_df.iloc[i]["VAC"]
        bac = labor_df.iloc[i]["BAC"]
        ratio = vac / bac if (not pd.isna(bac) and bac != 0) else np.nan
        rgb = vac_color_from_ratio(ratio)
        if rgb is not None:
            cell = labor_tbl.cell(i+1, vac_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb

    # Program Manpower table under labor table
    top_mp = top + labor_h + 0.2
    mp_label = slide2.shapes.add_textbox(
        Inches(0.5), Inches(top_mp - 0.3),
        Inches(3.0), Inches(0.3)
    )
    mp_tf = mp_label.text_frame
    mp_tf.text = "Program Manpower (FTE)"
    mp_p = mp_tf.paragraphs[0]
    mp_p.font.bold = True
    mp_p.font.size = Pt(12)

    mp_tbl, mp_h = add_simple_table(
        slide2, manpower_df,
        left_in=0.5, top_in=top_mp,
        width_in=6.0
    )

    # Color manpower % Var
    if "% Var" in manpower_df.columns:
        var_idx = list(manpower_df.columns).index("% Var")
        for i in range(len(manpower_df)):
            var_ratio = manpower_df.iloc[i]["% Var"]
            rgb = manpower_var_color(var_ratio)
            if rgb is not None:
                cell = mp_tbl.cell(i+1, var_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb

    # Comments box at bottom of slide 2
    bottom_top2 = top_mp + mp_h + 0.2
    add_rcca_box(slide2,
                 label="Comments / Root Cause & Corrective Actions",
                 left_in=0.5, top_in=bottom_top2, width_in=9.0, height_in=1.0)

    # ---------------------------------------------------
    # Slide 3 – Sub Team CPI/SPI/BEI Metrics
    # ---------------------------------------------------
    slide3 = prs.slides.add_slide(blank_layout)
    add_title(slide3, f"{program} EVMS Detail – Sub Team CPI / SPI / BEI Metrics")

    top3 = 0.9
    metrics_tbl, metrics_h = add_simple_table(
        slide3, metrics_sub,
        left_in=0.5, top_in=top3,
        width_in=9.0
    )

    # Color CPI/SPI/BEI cells
    cols_m = list(metrics_sub.columns)
    metric_cols = ["CPI CTD","CPI LSD","SPI CTD","SPI LSD","BEI CTD","BEI LSD"]
    metric_indices = {c: cols_m.index(c) for c in metric_cols}
    for i in range(len(metrics_sub)):
        for name, j in metric_indices.items():
            val = metrics_sub.iloc[i][name]
            rgb = spi_cpi_color(val)
            if rgb is not None:
                cell = metrics_tbl.cell(i+1, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb

    # Comments box on slide 3
    bottom_top3 = top3 + metrics_h + 0.2
    add_rcca_box(slide3,
                 label="Comments / Root Cause & Corrective Actions",
                 left_in=0.5, top_in=bottom_top3, width_in=9.0, height_in=1.0)

    # Save deck
    out_pptx = os.path.join(OUTPUT_DIR, f"{program}_EVMS_Deck.pptx")
    prs.save(out_pptx)
    print(f"  → Saved deck: {out_pptx}")

print("ALL PROGRAM EVMS DECKS COMPLETE ✓")